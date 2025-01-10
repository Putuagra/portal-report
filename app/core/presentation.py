from pptx import Presentation
import pandas as pd
from pptx.chart.data import CategoryChartData
from io import BytesIO
import matplotlib.pyplot as plt
import numpy as np

# Read the presentation
def read_presentation(file_path: str) -> Presentation:
    return Presentation(file_path)

# Get the slides
def get_slides(prs: Presentation) -> list:
    slides = prs.slides
    return slides

# Change the title of the slide
def change_title(slide: Presentation, new_title: str) -> Presentation:
    title = slide.shapes.title
    title.text = new_title
    return slide

# Find all charts
def find_charts(slide: Presentation) -> list:
    charts = []
    for shape in slide.shapes:
        if shape.has_chart:
            charts.append(shape)
    return charts

# Insert data into a chart
def insert_data(shape, dataframe, categories_column, series_name_column, values_column):
    """
    Insert data into a PowerPoint chart from a Pandas DataFrame.

    Parameters:
    - chart: The chart object to modify.
    - dataframe: Pandas DataFrame containing the data to populate the chart.
    - categories_column: Column name for the categories (x-axis labels).
    - series_name_column: Column name for the series names.
    - values_column: Column name for the values (y-axis data).

    Returns:
    - chart: The modified chart object.
    """

    chart = shape.chart

    # Create new chart data
    chart_data = CategoryChartData()

    # Add series data
    unique_series = dataframe[series_name_column].unique()
    for series in unique_series:
        series_data = dataframe[dataframe[series_name_column] == series]
        chart_data.add_series(series, series_data[values_column].tolist())

    # Add categories
    chart_data.categories = dataframe[categories_column].unique().tolist()

    # Replace the chart's data
    chart.replace_data(chart_data)

    return chart

# Run the presentation creation process
def run_presentation(
    transaction_data: pd.DataFrame,
    engagement_data: pd.DataFrame,
    conversion_data: pd.DataFrame,
) -> None:
    """
    Generates a presentation using transaction, engagement, and conversion data.

    Parameters:
    - prs: A PowerPoint Presentation object.
    - transaction_data: DataFrame containing transaction-related data.
    - engagement_data: DataFrame containing engagement-related data.
    - conversion_data: DataFrame containing conversion-related data.

    Returns:
    - None. Saves the generated presentation as 'output.pptx'.
    """

    # Load the presentation template
    path = "assets/presentation-template.pptx"  
    presentation = read_presentation(path)  

    # Retrieve all slides from the presentation
    slides = get_slides(presentation)  
    second_slide = slides[1] 

    # Locate the three charts on the second slide
    chart_one, chart_two, chart_three = find_charts(second_slide)  

    # Extract unique company name and year from the transaction data
    company_name = transaction_data['CompanyName'].unique()[0]
    year = engagement_data['Year'].unique()[0]

    # Update the slide title with the company name and year
    change_title(second_slide, f"{company_name} - {year}")  

    # Populate the charts 
    chart_one = insert_data(
        chart_one, transaction_data, "ServiceUsed", "ServiceUsed", "TransactionAmount"
    )

    chart_two = insert_data(
        chart_two, engagement_data, "TransactionDate", "ServiceUsed", "EngagementRate"
    )

    chart_three = insert_data(
        chart_three, conversion_data, "TransactionDate", "ServiceUsed", "ConversionRate"
    )
    
    ppt_buffer = BytesIO()
    presentation.save(ppt_buffer)
    ppt_buffer.seek(0)

    return ppt_buffer

def generate_vertical_bar(df, filename, label):
    df = df.head(20)
    plt.figure(figsize=(10, 6))
    width = 0.2
    x = np.arange(len(df["Transaction Date"]))
    if label == "Max TPS":
        bar1 = plt.bar(x - width, df[label], width=width, color="skyblue", label=label)
        bar2 = plt.bar(x, df["Avg TPS"], width=width, color="yellow", label="Avg TPS")
        bar3 = plt.bar(
            x + width,
            df["Max TPS (95th Percentile)"],
            width=width,
            color="red",
            label="Max TPS (95th Percentile)",
        )
        plt.title("TPS per Day")
        for bars in [bar1, bar2, bar3]:
            for bar in bars:
                height = bar.get_height()
                plt.text(
                    bar.get_x() + bar.get_width() / 2,
                    height,
                    f"{height:.0f}",
                    ha="center",
                    va="bottom",
                    fontsize=9,
                )
    else:
        bar = plt.bar(x, df[label], width=width, color="skyblue", label=label)
        plt.title(label)
        for bar in bar:
            height = bar.get_height()
            plt.text(
                bar.get_x() + bar.get_width() / 2,
                height,
                f"{height:.0f}",
                ha="center",
                va="bottom",
                fontsize=9,
            )

    # Add labels and title
    plt.xlabel("Transaction Date")
    plt.legend(bbox_to_anchor=(1, 1), loc="upper left")
    plt.xticks(x, df["Transaction Date"], rotation=45)
    plt.tight_layout()
    plt.savefig(filename, dpi=300, bbox_inches="tight", pad_inches=0)

    # Show the plot
    # plt.show()